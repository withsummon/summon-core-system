# Copyright (c) 2023-present Plane Software, Inc. and contributors
# SPDX-License-Identifier: AGPL-3.0-only
# See the LICENSE file for details.

import re
from io import BytesIO

import pytest

from plane.summon.services.document_renderer import render_document_files


pytestmark = [
    pytest.mark.unit,
    pytest.mark.filterwarnings("ignore:datetime.datetime.utcnow.*:DeprecationWarning:openpyxl.*"),
]

MARKDOWN = """# Delivery Plan

Prepared for the client.

## Scope

- Discovery
- Implementation

| Workstream | Owner |
| --- | --- |
| Platform | Summon |
| Security | IGLO |
"""


@pytest.mark.parametrize(
    ("document_type", "expected_formats"),
    [
        ("mom_iglo", ["docx", "pdf"]),
        ("mom_summon", ["docx", "pdf"]),
        ("proposal_vendor", ["docx", "pdf"]),
        ("proposal_client", ["docx", "pdf"]),
        ("invoice", ["docx", "pdf"]),
        ("quotation", ["docx", "pdf"]),
        ("uat", ["docx", "pdf"]),
        ("bast", ["docx", "pdf"]),
        ("presentation", ["pptx", "pdf"]),
        ("usage_cost", ["xlsx", "pdf"]),
        ("cost_projection", ["xlsx", "pdf"]),
        ("timeline", ["xlsx", "pdf"]),
        ("bug_report", ["xlsx", "pdf"]),
    ],
)
def test_document_types_render_the_approved_file_pairs(document_type, expected_formats):
    rendered = render_document_files(document_type, "Quarterly / Report", MARKDOWN)

    assert [item.format for item in rendered] == expected_formats
    assert [item.filename for item in rendered] == [f"quarterly-report.{format}" for format in expected_formats]
    assert all(item.data for item in rendered)


def test_docx_preserves_headings_paragraphs_bullets_tables_and_iglo_theme():
    from docx import Document
    from docx.oxml.ns import qn

    rendered = render_document_files("mom_iglo", "Weekly Meeting", MARKDOWN)[0]
    document = Document(BytesIO(rendered.data))
    paragraphs = {paragraph.text: paragraph for paragraph in document.paragraphs}

    assert rendered.content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    assert "Delivery Plan" in paragraphs
    assert "Prepared for the client." in paragraphs
    assert paragraphs["Discovery"].style.name == "List Bullet"
    assert [[cell.text for cell in row.cells] for row in document.tables[0].rows] == [
        ["Workstream", "Owner"],
        ["Platform", "Summon"],
        ["Security", "IGLO"],
    ]
    assert document.sections[0].header.paragraphs[0].text == "IGLO"
    assert str(document.paragraphs[0].runs[0].font.color.rgb) == "7A1F2B"
    shading = document.tables[0].rows[0].cells[0]._tc.get_or_add_tcPr().find(qn("w:shd"))
    assert shading.get(qn("w:fill")) == "F4C542"


def test_presentation_is_widescreen_and_contains_markdown_content():
    from pptx import Presentation

    rendered = render_document_files("presentation", "Delivery Review", MARKDOWN)[0]
    presentation = Presentation(BytesIO(rendered.data))
    text = "\n".join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text"))
    tables = [shape.table for slide in presentation.slides for shape in slide.shapes if shape.has_table]

    assert rendered.content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    assert presentation.slide_width / presentation.slide_height == pytest.approx(16 / 9)
    assert "Delivery Review" in text
    assert "Scope" in text
    assert "Discovery" in text
    assert [[cell.text for cell in row.cells] for row in tables[0].rows] == [
        ["Workstream", "Owner"],
        ["Platform", "Summon"],
        ["Security", "IGLO"],
    ]
    assert "2026 • CONFIDENTIAL" in text
    assert "withsummon.com" in text


def test_spreadsheet_preserves_blocks_and_neutralizes_formula_cells():
    from openpyxl import load_workbook

    markdown = """# Usage

- Metered usage

| Item | Cost |
| --- | --- |
| API | =SUM(1,1) |
"""
    rendered = render_document_files("usage_cost", "=Usage Cost", markdown)[0]
    worksheet = load_workbook(BytesIO(rendered.data)).active
    values = [cell.value for row in worksheet.iter_rows() for cell in row if cell.value is not None]

    assert rendered.content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    assert "'=Usage Cost" in values
    assert "Metered usage" in values
    assert "Item" in values
    assert "'=SUM(1,1)" in values
    assert all(cell.data_type != "f" for row in worksheet.iter_rows() for cell in row)


def test_pdf_uses_a4_and_unknown_document_types_are_rejected():
    rendered = render_document_files("proposal_client", "Proposal", MARKDOWN)[1]
    media_box = re.search(rb"/MediaBox\s*\[\s*0\s+0\s+([\d.]+)\s+([\d.]+)\s*\]", rendered.data)

    assert rendered.filename == "proposal.pdf"
    assert rendered.content_type == "application/pdf"
    assert rendered.data.startswith(b"%PDF-")
    assert media_box
    assert float(media_box.group(1)) == pytest.approx(595.28, abs=0.02)
    assert float(media_box.group(2)) == pytest.approx(841.89, abs=0.02)

    with pytest.raises(ValueError, match="Unsupported document type"):
        render_document_files("unknown", "Unsafe", MARKDOWN)


def test_presentation_pdf_is_multipage_widescreen():
    rendered = render_document_files("presentation", "Delivery Review", MARKDOWN)[1]
    pages = re.findall(rb"/Type\s*/Page\b", rendered.data)
    media_boxes = re.findall(rb"/MediaBox\s*\[\s*0\s+0\s+([\d.]+)\s+([\d.]+)\s*\]", rendered.data)

    assert len(pages) > 1
    assert len(media_boxes) == len(pages)
    assert all(float(width) / float(height) == pytest.approx(16 / 9) for width, height in media_boxes)


def test_markdown_preserves_literal_symbols_escapes_and_fenced_code():
    from docx import Document

    markdown = r"""# Symbols

SKU_A_B costs 2 * 3; trace_id is `trace_id`; keep a lone ~ and ` tick.

**Bold** *italic* ~~removed~~ [Portal](https://withsummon.com)

\*literal\* \_escaped\_ \`tick\`

```text
# not a heading
SKU_A_B * raw `trace_id`
| not | a table |
```
"""
    rendered = render_document_files("mom_summon", "Markdown Safety", markdown)[0]
    document = Document(BytesIO(rendered.data))
    paragraphs = [paragraph.text for paragraph in document.paragraphs]

    assert "SKU_A_B costs 2 * 3; trace_id is trace_id; keep a lone ~ and ` tick." in paragraphs
    assert "Bold italic removed Portal" in paragraphs
    assert "*literal* _escaped_ `tick`" in paragraphs
    assert "# not a heading\nSKU_A_B * raw `trace_id`\n| not | a table |" in paragraphs
    assert all(
        paragraph.style.name != "Heading 1" for paragraph in document.paragraphs if "not a heading" in paragraph.text
    )


def test_long_presentation_creates_bounded_continuation_slides_and_pdf_pages():
    from pptx import Presentation
    from pptx.util import Inches

    bullets = "\n".join(
        f"- Delivery workstream {index}: owner confirms dependencies, evidence, risks, and acceptance criteria."
        for index in range(1, 37)
    )
    rows = "\n".join(
        f"| Phase {index} | Owner {index} | Evidence and final acceptance marker {index} |" for index in range(1, 25)
    )
    markdown = f"""# Delivery Plan

{bullets}

| Phase | Owner | Evidence |
| --- | --- | --- |
{rows}
"""
    pptx, pdf = render_document_files("presentation", "Long Delivery Review", markdown)
    presentation = Presentation(BytesIO(pptx.data))
    pages = re.findall(rb"/Type\s*/Page\b", pdf.data)
    table_text = "\n".join(
        cell.text
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_table
        for row in shape.table.rows
        for cell in row.cells
    )

    assert len(presentation.slides) > 3
    assert len(pages) == len(presentation.slides)
    assert "final acceptance marker 24" in table_text
    for slide in list(presentation.slides)[1:]:
        content_shapes = [
            shape
            for shape in slide.shapes
            if not (hasattr(shape, "text") and shape.text in {"withsummon.com", "2026 • CONFIDENTIAL"})
        ]
        assert all(shape.top + shape.height <= Inches(6.6) for shape in content_shapes)


@pytest.mark.parametrize(
    ("document_type", "label", "section", "columns"),
    [
        ("mom_iglo", "IGLO Minutes of Meeting", "Discussion", ("No", "Tugas", "Keterangan")),
        ("mom_summon", "Summon Minutes of Meeting", "Decisions", ("No", "Tugas", "Keterangan")),
        ("proposal_vendor", "Vendor Proposal", "Executive Summary", ("Phase", "Deliverable", "Timeline")),
        ("proposal_client", "Client Proposal", "Pricing Scheme", ("Item", "Description", "Amount")),
        ("invoice", "Invoice", "Bill To", ("Description", "Qty", "Unit Rate", "Amount")),
        ("quotation", "Quotation", "Terms", ("Item", "Description", "Amount")),
        (
            "uat",
            "User Acceptance Test",
            "Changes Being Tested",
            ("#", "Test Case", "Steps", "Expected Result", "Status", "Notes"),
        ),
        ("bast", "Handover and Acceptance (BAST)", "Deliverables", ("Deliverable", "Progress", "Status")),
        ("presentation", "Summon Presentation", "Next Steps", ("Slide", "Key Message", "Evidence")),
        (
            "usage_cost",
            "Usage Cost",
            "Period Summary",
            ("Timestamp", "Session", "Source", "Minutes", "Input Tokens", "Output Tokens", "Total USD"),
        ),
        (
            "cost_projection",
            "Cost Projection",
            "Monthly Scenarios",
            ("Scenario", "Sessions", "Cost USD", "Cost IDR"),
        ),
        ("timeline", "Project Timeline", "Milestones", ("No", "Scope of Work", "Week/Month", "Progress")),
        (
            "bug_report",
            "Bug Report",
            "Client Report",
            ("Date Reported", "What's Happening?", "Steps", "Environment", "App Version", "Urgency", "Status"),
        ),
    ],
)
def test_document_types_apply_their_canonical_layout_signature(document_type, label, section, columns):
    from docx import Document
    from openpyxl import load_workbook
    from pptx import Presentation

    header = " | ".join(columns)
    separator = " | ".join("---" for _ in columns)
    values = " | ".join(f"Value {index}" for index in range(1, len(columns) + 1))
    markdown = f"# {section}\n\nReference content.\n\n| {header} |\n| {separator} |\n| {values} |"
    rendered = render_document_files(document_type, "Reference Output", markdown)[0]

    if rendered.format == "docx":
        document = Document(BytesIO(rendered.data))
        assert label in "\n".join(paragraph.text for paragraph in document.sections[0].header.paragraphs)
        assert section in {paragraph.text for paragraph in document.paragraphs}
        assert [cell.text for cell in document.tables[0].rows[0].cells] == list(columns)
    elif rendered.format == "pptx":
        presentation = Presentation(BytesIO(rendered.data))
        text = "\n".join(
            shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")
        )
        table = next(shape.table for slide in presentation.slides for shape in slide.shapes if shape.has_table)
        assert label in text
        assert section in text
        assert [cell.text for cell in table.rows[0].cells] == list(columns)
    else:
        worksheet = load_workbook(BytesIO(rendered.data)).active
        contents = {cell.value for row in worksheet.iter_rows() for cell in row if cell.value is not None}
        assert worksheet.title == label
        assert section in contents
        assert set(columns) <= contents


def test_long_invoice_repeats_table_header_and_flows_to_multiple_pdf_pages():
    from docx import Document
    from docx.oxml.ns import qn

    rows = "\n".join(
        "| Managed service line "
        f"{index} with implementation and support | {index} | IDR 1,000,000 | IDR {index},000,000 |"
        for index in range(1, 46)
    )
    markdown = f"""# Bill To

PT Example Client, Jakarta

| Description | Qty | Unit Rate | Amount |
| --- | --- | --- | --- |
{rows}

## Payment Details

Payment follows the approved commercial terms. Final invoice acceptance marker.
"""
    docx, pdf = render_document_files("invoice", "INV-2026-0042", markdown)
    document = Document(BytesIO(docx.data))
    table = document.tables[0]
    header_properties = table.rows[0]._tr.get_or_add_trPr()

    assert len(table.rows) == 46
    assert header_properties.find(qn("w:tblHeader")) is not None
    assert len(re.findall(rb"/Type\s*/Page\b", pdf.data)) > 1
    assert any("Final invoice acceptance marker" in paragraph.text for paragraph in document.paragraphs)


@pytest.mark.parametrize(
    ("document_type", "columns"),
    [
        ("timeline", ("No", "Scope of Work", "Week/Month", "Progress")),
        (
            "bug_report",
            ("Date Reported", "What's Happening?", "Steps", "Environment", "App Version", "Urgency", "Status"),
        ),
    ],
)
def test_long_tabular_reports_use_bounded_landscape_print_layout(document_type, columns):
    from openpyxl import load_workbook

    header = " | ".join(columns)
    separator = " | ".join("---" for _ in columns)
    rows = "\n".join(
        "| " + " | ".join(f"Row {row} detailed acceptance evidence for {column}" for column in columns) + " |"
        for row in range(1, 41)
    )
    markdown = f"# Delivery Detail\n\n| {header} |\n| {separator} |\n{rows}"
    xlsx, pdf = render_document_files(document_type, "Long Delivery Detail", markdown)
    worksheet = load_workbook(BytesIO(xlsx.data)).active
    media_box = re.search(rb"/MediaBox\s*\[\s*0\s+0\s+([\d.]+)\s+([\d.]+)\s*\]", pdf.data)

    assert worksheet.page_setup.orientation == "landscape"
    assert worksheet.page_setup.fitToWidth == 1
    assert worksheet.print_title_rows == "$3:$3"
    assert max(dimension.width for dimension in worksheet.column_dimensions.values()) <= 24
    assert all(cell.alignment.wrap_text for row in worksheet.iter_rows() for cell in row if cell.value is not None)
    assert media_box and float(media_box.group(1)) > float(media_box.group(2))
    assert len(re.findall(rb"/Type\s*/Page\b", pdf.data)) > 1
