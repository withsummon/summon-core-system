# Copyright (c) 2023-present Plane Software, Inc. and contributors
# SPDX-License-Identifier: AGPL-3.0-only
# See the LICENSE file for details.

import re
import textwrap
import unicodedata
from dataclasses import dataclass
from functools import partial
from io import BytesIO
from xml.sax.saxutils import escape

from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Mm, Pt, RGBColor
from openpyxl import Workbook
from openpyxl.formatting.rule import FormulaRule
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.datavalidation import DataValidation
from pptx import Presentation
from pptx.dml.color import RGBColor as PresentationRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt as PresentationPt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch, mm
from reportlab.pdfgen.canvas import Canvas
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle


@dataclass(frozen=True)
class RenderedDocument:
    filename: str
    content_type: str
    data: bytes
    format: str


@dataclass(frozen=True)
class _Theme:
    primary: str
    accent: str
    header_text: str


@dataclass(frozen=True)
class _Profile:
    label: str
    landscape: bool = False


@dataclass(frozen=True)
class _Block:
    kind: str
    text: str = ""
    level: int = 0
    rows: tuple[tuple[str, ...], ...] = ()


SUMMON_THEME = _Theme("102A43", "2563EB", "FFFFFF")
IGLO_THEME = _Theme("7A1F2B", "F4C542", "7A1F2B")

DOCUMENT_FORMATS = {
    "mom_iglo": ("docx", "pdf"),
    "mom_summon": ("docx", "pdf"),
    "proposal_vendor": ("pptx", "pdf"),
    "proposal_client": ("pptx", "pdf"),
    "invoice": ("docx", "pdf"),
    "quotation": ("docx", "pdf"),
    "uat": ("docx", "pdf"),
    "bast": ("docx", "pdf"),
    "presentation": ("pptx", "pdf"),
    "usage_cost": ("xlsx", "pdf"),
    "cost_projection": ("xlsx", "pdf"),
    "timeline": ("xlsx", "pdf"),
    "bug_report": ("xlsx", "pdf"),
}

CONTENT_TYPES = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pdf": "application/pdf",
}

PROFILES = {
    "mom_iglo": _Profile("IGLO Minutes of Meeting"),
    "mom_summon": _Profile("Summon Minutes of Meeting"),
    "proposal_vendor": _Profile("Vendor Proposal"),
    "proposal_client": _Profile("Client Proposal"),
    "invoice": _Profile("Invoice"),
    "quotation": _Profile("Quotation"),
    "uat": _Profile("User Acceptance Test"),
    "bast": _Profile("Handover and Acceptance (BAST)"),
    "presentation": _Profile("Summon Presentation", True),
    "usage_cost": _Profile("Usage Cost", True),
    "cost_projection": _Profile("Cost Projection", True),
    "timeline": _Profile("Project Timeline", True),
    "bug_report": _Profile("Bug Report", True),
}


def _plain(text):
    escaped = []

    def protect(match):
        escaped.append(match.group(1))
        return f"\ue000{len(escaped) - 1}\ue001"

    text = re.sub(r"\\([\\`*_[\]{}()#+.!|>~-])", protect, text)
    text = re.sub(r"!\[([^]]*)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"\[([^]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"(?<!`)`([^`\n]+)`(?!`)", r"\1", text)
    text = re.sub(r"(?<!\*)\*\*(?=\S)(.+?\S)\*\*(?!\*)", r"\1", text)
    text = re.sub(r"(?<![\w_])__(?=\S)(.+?\S)__(?![\w_])", r"\1", text)
    text = re.sub(r"(?<!~)~~(?=\S)(.+?\S)~~(?!~)", r"\1", text)
    text = re.sub(r"(?<![\w*])\*(?=\S)([^*\n]*?\S)\*(?![\w*])", r"\1", text)
    text = re.sub(r"(?<![\w_])_(?=\S)([^_\n]*?\S)_(?![\w_])", r"\1", text)
    for index, value in enumerate(escaped):
        text = text.replace(f"\ue000{index}\ue001", value)
    return text.strip()


def _pipe_cells(line):
    stripped = line.strip()
    if "|" not in stripped:
        return None
    cells = re.split(r"(?<!\\)\|", stripped.strip("|"))
    return tuple(_plain(cell.replace(r"\|", "|")) for cell in cells)


def _is_table_separator(cells):
    return bool(cells) and all(re.fullmatch(r":?-{3,}:?", cell.replace(" ", "")) for cell in cells)


def _parse_markdown(markdown):
    lines = markdown.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    blocks = []
    paragraph = []

    def flush_paragraph():
        if paragraph:
            blocks.append(_Block("paragraph", " ".join(paragraph)))
            paragraph.clear()

    index = 0
    while index < len(lines):
        line = lines[index]
        fence = re.match(r"^\s*(`{3,}|~{3,})", line)
        if fence:
            flush_paragraph()
            marker = fence.group(1)
            code = []
            index += 1
            while index < len(lines) and not re.match(
                rf"^\s*{re.escape(marker[0])}{{{len(marker)},}}\s*$", lines[index]
            ):
                code.append(lines[index])
                index += 1
            blocks.append(_Block("paragraph", "\n".join(code)))
            index += index < len(lines)
            continue
        heading = re.match(r"^(#{1,6})\s+(.+)$", line.strip())
        bullet = re.match(r"^\s*[-*+]\s+(.+)$", line)
        cells = _pipe_cells(line)
        separator = _pipe_cells(lines[index + 1]) if index + 1 < len(lines) else None
        if cells and separator and len(cells) == len(separator) and _is_table_separator(separator):
            flush_paragraph()
            rows = [cells]
            index += 2
            while index < len(lines):
                row = _pipe_cells(lines[index])
                if not row:
                    break
                rows.append(tuple((*row, *("" for _ in range(max(0, len(cells) - len(row))))))[: len(cells)])
                index += 1
            blocks.append(_Block("table", rows=tuple(rows)))
            continue
        if not line.strip():
            flush_paragraph()
        elif heading:
            flush_paragraph()
            blocks.append(_Block("heading", _plain(heading.group(2)), len(heading.group(1))))
        elif bullet:
            flush_paragraph()
            blocks.append(_Block("bullet", _plain(bullet.group(1))))
        else:
            paragraph.append(_plain(line))
        index += 1
    flush_paragraph()
    return blocks


def _filename(title):
    ascii_title = unicodedata.normalize("NFKD", title).encode("ascii", "ignore").decode().lower()
    return re.sub(r"[^a-z0-9]+", "-", ascii_title).strip("-") or "document"


def _docx(title, blocks, theme, profile):
    document = Document()
    section = document.sections[0]
    section.page_width, section.page_height = Mm(210), Mm(297)
    section.left_margin = section.right_margin = Mm(22)
    if theme == IGLO_THEME:
        header_run = section.header.paragraphs[0].add_run("IGLO")
        header_run.bold = True
        header_run.font.color.rgb = RGBColor.from_string(theme.primary)
        signature = section.header.add_paragraph()
    else:
        signature = section.header.paragraphs[0]
    signature_run = signature.add_run(profile.label)
    signature_run.bold = True
    signature_run.font.color.rgb = RGBColor.from_string(theme.primary)
    document.core_properties.title = title
    title_paragraph = document.add_paragraph(style="Title")
    title_run = title_paragraph.add_run(title)
    title_run.bold = True
    title_run.font.size = Pt(26)
    title_run.font.color.rgb = RGBColor.from_string(theme.primary)
    for block in blocks:
        if block.kind == "heading":
            paragraph = document.add_heading(block.text, level=min(block.level, 3))
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor.from_string(theme.primary)
        elif block.kind == "paragraph":
            document.add_paragraph(block.text)
        elif block.kind == "bullet":
            document.add_paragraph(block.text, style="List Bullet")
        else:
            table = document.add_table(rows=len(block.rows), cols=len(block.rows[0]))
            table.style = "Table Grid"
            header_properties = table.rows[0]._tr.get_or_add_trPr()
            repeat_header = OxmlElement("w:tblHeader")
            repeat_header.set(qn("w:val"), "true")
            header_properties.append(repeat_header)
            for row_index, row in enumerate(block.rows):
                for column_index, value in enumerate(row):
                    cell = table.cell(row_index, column_index)
                    cell.text = value
                    if row_index == 0:
                        shading = OxmlElement("w:shd")
                        shading.set(qn("w:fill"), theme.accent)
                        cell._tc.get_or_add_tcPr().append(shading)
                        for run in cell.paragraphs[0].runs:
                            run.bold = True
                            run.font.color.rgb = RGBColor.from_string(theme.header_text)
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def _sections(title, blocks):
    sections = []
    heading, content = title, []
    for block in blocks:
        if block.kind == "heading":
            if content or heading != title:
                sections.append((heading, content))
            heading, content = block.text, []
        else:
            content.append(block)
    if content or not sections:
        sections.append((heading, content))
    return sections


SLIDE_CONTENT_HEIGHT = 5.0


def _wrapped_lines(text, width=88):
    lines = []
    for line in text.splitlines() or ("",):
        lines.extend(textwrap.wrap(line, width=width, break_long_words=True, break_on_hyphens=False) or [""])
    return lines


def _table_row_height(row):
    cell_width = max(18, 100 // len(row))
    lines = max(len(_wrapped_lines(value, cell_width)) for value in row)
    return max(0.38, 0.32 * lines + 0.14)


def _block_height(block):
    if block.kind == "table":
        return sum(_table_row_height(row) for row in block.rows) + 0.12
    return max(0.36, 0.32 * len(_wrapped_lines(block.text)) + 0.08)


def _slide_block_chunks(block):
    if block.kind == "table":
        header, *rows = block.rows
        chunks, current = [], [header]
        height = _table_row_height(header) + 0.12
        for row in rows:
            row_height = _table_row_height(row)
            if len(current) > 1 and height + row_height > SLIDE_CONTENT_HEIGHT:
                chunks.append(_Block("table", rows=tuple(current)))
                current, height = [header], _table_row_height(header) + 0.12
            current.append(row)
            height += row_height
        if len(current) > 1 or not chunks:
            chunks.append(_Block("table", rows=tuple(current)))
        return chunks
    lines = _wrapped_lines(block.text)
    lines_per_slide = max(1, int((SLIDE_CONTENT_HEIGHT - 0.08) / 0.32))
    return [
        _Block(block.kind, "\n".join(lines[index : index + lines_per_slide]), block.level)
        for index in range(0, len(lines), lines_per_slide)
    ]


def _slide_pages(title, blocks):
    pages = []
    for heading, content in _sections(title, blocks):
        current, height, page_number = [], 0.0, 1
        for block in content:
            for chunk in _slide_block_chunks(block):
                chunk_height = _block_height(chunk)
                if current and height + chunk_height > SLIDE_CONTENT_HEIGHT:
                    pages.append((heading if page_number == 1 else f"{heading} (continued)", tuple(current)))
                    current, height, page_number = [], 0.0, page_number + 1
                current.append(chunk)
                height += chunk_height
        pages.append((heading if page_number == 1 else f"{heading} (continued)", tuple(current)))
    return pages


def _pptx(title, blocks, theme, profile):
    presentation = Presentation()
    slide_width, slide_height = (13.333333, 7.5) if profile.landscape else (8.267717, 11.692913)
    content_width = slide_width - 1.6
    presentation.slide_width, presentation.slide_height = Inches(slide_width), Inches(slide_height)
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = title
    title_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PresentationRGBColor.from_string(theme.primary)
    title_slide.placeholders[1].text = profile.label
    for heading, content in _slide_pages(title, blocks):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(slide_width - 1.4), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = heading
        title_frame.paragraphs[0].font.size = PresentationPt(26)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = PresentationRGBColor.from_string(theme.primary)
        y = 1.35
        for block in content:
            block_height = _block_height(block)
            if block.kind == "table":
                rows, columns = len(block.rows), len(block.rows[0])
                shape = slide.shapes.add_table(
                    rows, columns, Inches(0.8), Inches(y), Inches(content_width), Inches(block_height - 0.12)
                )
                for row_index, row in enumerate(block.rows):
                    shape.table.rows[row_index].height = Inches(_table_row_height(row))
                    for column_index, value in enumerate(row):
                        cell = shape.table.cell(row_index, column_index)
                        cell.text = value
                        cell.text_frame.paragraphs[0].font.size = PresentationPt(14)
                        if row_index == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = PresentationRGBColor.from_string(theme.accent)
                            cell.text_frame.paragraphs[0].font.bold = True
                            cell.text_frame.paragraphs[0].font.color.rgb = PresentationRGBColor.from_string(
                                theme.header_text
                            )
            else:
                text_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(y), Inches(content_width), Inches(block_height - 0.08)
                )
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                text_frame.margin_left = text_frame.margin_right = 0
                text_frame.margin_top = text_frame.margin_bottom = 0
                paragraph = text_frame.paragraphs[0]
                paragraph.text = f"• {block.text}" if block.kind == "bullet" else block.text
                paragraph.font.size = PresentationPt(18)
                paragraph.font.color.rgb = PresentationRGBColor.from_string(SUMMON_THEME.primary)
            y += block_height
    for slide in presentation.slides:
        for text, left, alignment in (
            ("withsummon.com", 0.8, PP_ALIGN.LEFT),
            ("2026 • CONFIDENTIAL", slide_width - 3.9, PP_ALIGN.RIGHT),
        ):
            footer = slide.shapes.add_textbox(
                Inches(left), Inches(slide_height - 0.45), Inches(3.1), Inches(0.2)
            ).text_frame
            footer.text = text
            footer.paragraphs[0].alignment = alignment
            footer.paragraphs[0].font.size = PresentationPt(9)
            footer.paragraphs[0].font.color.rgb = PresentationRGBColor.from_string(theme.primary)
    if not profile.landscape:
        for slide in presentation.slides:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = PresentationRGBColor.from_string(theme.primary)
            for shape in slide.shapes:
                if shape.has_table or not hasattr(shape, "text_frame"):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = PresentationRGBColor(255, 255, 255)
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _spreadsheet_value(value):
    return f"'{value}" if value.startswith(("=", "+", "-", "@", "\t", "\r", "\n")) else value


def _xlsx(title, blocks, theme, profile):
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = profile.label[:31]
    worksheet.page_setup.orientation = worksheet.ORIENTATION_LANDSCAPE
    worksheet.page_setup.paperSize = worksheet.PAPERSIZE_A4
    worksheet.page_setup.fitToWidth = 1
    worksheet.page_setup.fitToHeight = 0
    worksheet.sheet_properties.pageSetUpPr.fitToPage = True
    worksheet.sheet_view.showGridLines = False
    width = max((len(row) for block in blocks if block.kind == "table" for row in block.rows), default=2)
    worksheet.append([_spreadsheet_value(title)])
    if width > 1:
        worksheet.merge_cells(start_row=1, start_column=1, end_row=1, end_column=width)
    worksheet["A1"].font = Font(size=18, bold=True, color="FFFFFF")
    worksheet["A1"].fill = PatternFill("solid", fgColor=theme.primary)
    worksheet["A1"].alignment = Alignment(vertical="center")
    table_header_row = None
    for block in blocks:
        if block.kind in {"heading", "paragraph"}:
            worksheet.append([_spreadsheet_value(block.text)])
            row = worksheet.max_row
            if width > 1:
                worksheet.merge_cells(start_row=row, start_column=1, end_row=row, end_column=width)
            if block.kind == "heading":
                worksheet.cell(row, 1).font = Font(bold=True, color="FFFFFF")
                worksheet.cell(row, 1).fill = PatternFill("solid", fgColor=theme.primary)
        elif block.kind == "bullet":
            worksheet.append(["•", _spreadsheet_value(block.text)])
        else:
            for row_index, values in enumerate(block.rows):
                worksheet.append([_spreadsheet_value(value) for value in values])
                if row_index == 0:
                    table_header_row = table_header_row or worksheet.max_row
                    for cell in worksheet[worksheet.max_row]:
                        cell.font = Font(bold=True, color=theme.header_text)
                        cell.fill = PatternFill("solid", fgColor=theme.accent)
    if table_header_row:
        worksheet.print_title_rows = f"{table_header_row}:{table_header_row}"
        worksheet.freeze_panes = f"A{table_header_row + 1}"
    for row in worksheet.iter_rows():
        for cell in row:
            if cell.value is not None:
                cell.alignment = Alignment(vertical="top", wrap_text=True)
    for column in range(1, width + 1):
        longest = max(
            (
                max((len(line) for line in str(cell.value).splitlines()), default=0)
                for cell in worksheet[get_column_letter(column)]
            ),
            default=0,
        )
        worksheet.column_dimensions[get_column_letter(column)].width = min(24, max(12, longest + 2))
    if profile.label == "Bug Report":
        worksheet.title = "Tracker"
        if table_header_row:
            last_column = get_column_letter(width)
            worksheet.auto_filter.ref = f"A{table_header_row}:{last_column}{worksheet.max_row}"
            headers = {str(worksheet.cell(table_header_row, column).value): column for column in range(1, width + 1)}
            validations = {
                "Severity": '"Low,Medium,High,Critical"',
                "Current Status": '"New,Triaged,In Progress,Ready for Retest,Closed,Blocked"',
                "Client Verification": '"Not Tested,Verified,Rejected"',
            }
            for header, formula in validations.items():
                column = headers.get(header)
                if not column:
                    continue
                validation = DataValidation(type="list", formula1=formula, allow_blank=True)
                worksheet.add_data_validation(validation)
                letter = get_column_letter(column)
                validation.add(f"{letter}{table_header_row + 1}:{letter}500")
            status_column = headers.get("Current Status")
            if status_column:
                letter = get_column_letter(status_column)
                cell_range = f"{letter}{table_header_row + 1}:{letter}500"
                for status, color in (("Closed", "DCFCE7"), ("Blocked", "FEE2E2"), ("In Progress", "FEF3C7")):
                    worksheet.conditional_formatting.add(
                        cell_range,
                        FormulaRule(
                            formula=[f'${letter}{table_header_row + 1}="{status}"'],
                            fill=PatternFill("solid", fgColor=color),
                        ),
                    )
        guide = workbook.create_sheet("How to Use", 0)
        guide.append(["Bug Report Workbook"])
        guide.append([])
        guide.append(
            ["Client-entered fields stay separate from developer-maintained investigation and delivery fields."]
        )
        guide.append(["Do not include credentials, session data, tokens, or unredacted secrets in evidence or notes."])
        guide.append(
            ["Use the dropdowns for severity, status, and client verification; leave unknown facts blank or TBD."]
        )
        guide.column_dimensions["A"].width = 110
        guide["A1"].font = Font(size=18, bold=True, color="FFFFFF")
        guide["A1"].fill = PatternFill("solid", fgColor=theme.primary)
        for cell in guide["A"]:
            cell.alignment = Alignment(vertical="top", wrap_text=True)
        workbook.active = workbook.index(worksheet)
    output = BytesIO()
    workbook.save(output)
    return output.getvalue()


def _pdf(title, blocks, theme, profile):
    output = BytesIO()
    document = SimpleDocTemplate(
        output,
        pagesize=landscape(A4) if profile.landscape else A4,
        leftMargin=18 * mm,
        rightMargin=18 * mm,
        topMargin=18 * mm,
        bottomMargin=18 * mm,
        pageCompression=0,
        title=title,
    )
    base = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "SummonTitle",
        parent=base["Title"],
        fontName="Helvetica-Bold",
        fontSize=22,
        textColor=colors.HexColor(f"#{theme.primary}"),
    )
    heading_style = ParagraphStyle(
        "SummonHeading",
        parent=base["Heading2"],
        fontName="Helvetica-Bold",
        textColor=colors.HexColor(f"#{theme.primary}"),
    )
    body_style = ParagraphStyle("SummonBody", parent=base["BodyText"], leading=15, spaceAfter=7)
    bullet_style = ParagraphStyle("SummonBullet", parent=body_style, leftIndent=12, firstLineIndent=-8)
    story = [Paragraph(escape(profile.label), heading_style), Paragraph(escape(title), title_style), Spacer(1, 4 * mm)]
    for block in blocks:
        if block.kind == "heading":
            story.append(Paragraph(escape(block.text), heading_style))
        elif block.kind == "paragraph":
            story.append(Paragraph(escape(block.text), body_style))
        elif block.kind == "bullet":
            story.append(Paragraph(escape(block.text), bullet_style, bulletText="•"))
        else:
            cells = [[Paragraph(escape(value), body_style) for value in row] for row in block.rows]
            table = Table(cells, colWidths=[document.width / len(block.rows[0])] * len(block.rows[0]), repeatRows=1)
            table.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(f"#{theme.accent}")),
                        ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor(f"#{theme.header_text}")),
                        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#CBD5E1")),
                        ("VALIGN", (0, 0), (-1, -1), "TOP"),
                        ("ALIGN", (0, 0), (-1, 0), "CENTER"),
                    ]
                )
            )
            story.extend([table, Spacer(1, 3 * mm)])
    document.build(story, canvasmaker=partial(Canvas, invariant=1))
    return output.getvalue()


def _presentation_pdf(title, blocks, theme, profile):
    output = BytesIO()
    page_size = (13.333333 * inch, 7.5 * inch) if profile.landscape else A4
    canvas = Canvas(output, pagesize=page_size, invariant=1, pageCompression=0, title=title)
    width, height = page_size
    content_width = width - 1.6 * inch
    primary = colors.HexColor(f"#{theme.primary}")
    foreground = colors.white if not profile.landscape else colors.HexColor("#102A43")
    body_style = ParagraphStyle("SlideBody", fontName="Helvetica", fontSize=14, leading=20, textColor=foreground)
    header_style = ParagraphStyle("SlideTableHeader", parent=body_style, textColor=colors.HexColor("#FFFFFF"))
    table_style = TableStyle(
        [
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(f"#{theme.accent}")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor(f"#{theme.header_text}")),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#CBD5E1")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ]
    )

    def start_page():
        if not profile.landscape:
            canvas.setFillColor(primary)
            canvas.rect(0, 0, width, height, fill=1, stroke=0)

    def finish_page():
        canvas.setFillColor(colors.white if not profile.landscape else primary)
        canvas.setFont("Helvetica", 9)
        canvas.drawString(0.8 * inch, 0.28 * inch, "withsummon.com")
        canvas.drawRightString(width - 0.8 * inch, 0.28 * inch, "2026 • CONFIDENTIAL")
        canvas.showPage()

    start_page()
    canvas.setFillColor(colors.white if not profile.landscape else primary)
    canvas.setFont("Helvetica-Bold", 30)
    canvas.drawString(0.8 * inch, height - 2.6 * inch, title)
    canvas.setFont("Helvetica", 16)
    canvas.drawString(0.8 * inch, height - 3 * inch, profile.label)
    finish_page()

    for heading, content in _slide_pages(title, blocks):
        start_page()
        canvas.setFillColor(colors.white if not profile.landscape else primary)
        canvas.setFont("Helvetica-Bold", 26)
        canvas.drawString(0.7 * inch, height - 0.85 * inch, heading)
        y = height - 1.55 * inch
        for block in content:
            if block.kind == "table":
                cells = [
                    [Paragraph(escape(value), header_style if row_index == 0 else body_style) for value in row]
                    for row_index, row in enumerate(block.rows)
                ]
                table = Table(cells, colWidths=[content_width / len(block.rows[0])] * len(block.rows[0]))
                table.setStyle(table_style)
                _, table_height = table.wrapOn(canvas, content_width, y)
                y -= table_height
                table.drawOn(canvas, 0.8 * inch, y)
                y -= 0.12 * inch
            else:
                prefix = "• " if block.kind == "bullet" else ""
                paragraph = Paragraph(f"{prefix}{escape(block.text).replace(chr(10), '<br/>')}", body_style)
                _, paragraph_height = paragraph.wrap(content_width, y)
                y -= paragraph_height
                paragraph.drawOn(canvas, 0.8 * inch, y)
                y -= 0.08 * inch
        finish_page()
    canvas.save()
    return output.getvalue()


def render_document_files(document_type: str, title: str, markdown: str) -> list[RenderedDocument]:
    normalized_type = document_type.strip().lower()
    formats = DOCUMENT_FORMATS.get(normalized_type)
    if not formats:
        raise ValueError(f"Unsupported document type: {document_type}")
    safe_title = title.strip() or "Document"
    blocks = _parse_markdown(markdown)
    theme = IGLO_THEME if normalized_type == "mom_iglo" else SUMMON_THEME
    profile = PROFILES[normalized_type]
    renderers = {"docx": _docx, "pptx": _pptx, "xlsx": _xlsx, "pdf": _pdf}
    if normalized_type in {"presentation", "proposal_vendor", "proposal_client"}:
        renderers["pdf"] = _presentation_pdf
    base_filename = _filename(safe_title)
    return [
        RenderedDocument(
            filename=f"{base_filename}.{format_name}",
            content_type=CONTENT_TYPES[format_name],
            data=renderers[format_name](safe_title, blocks, theme, profile),
            format=format_name,
        )
        for format_name in formats
    ]
