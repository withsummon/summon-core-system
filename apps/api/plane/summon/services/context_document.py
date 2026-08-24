from pathlib import Path
from zipfile import BadZipFile, ZipFile

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from rest_framework import serializers

MAX_UPLOAD_BYTES = 10 * 1024 * 1024
MAX_EXPANDED_BYTES = 50 * 1024 * 1024
MAX_CONTEXT_CHARS = 30000
SUPPORTED_EXTENSIONS = {".csv", ".docx", ".md", ".pdf", ".pptx", ".txt", ".xlsx"}


def _bounded_text(lines):
    text = ""
    for line in lines:
        line = str(line or "").strip()
        if not line:
            continue
        addition = ("\n" if text else "") + line
        if len(text) + len(addition) > MAX_CONTEXT_CHARS:
            return (text + addition)[:MAX_CONTEXT_CHARS], True
        text += addition
    return text, False


def _validate_office_archive(upload):
    try:
        upload.seek(0)
        with ZipFile(upload) as archive:
            if sum(item.file_size for item in archive.infolist()) > MAX_EXPANDED_BYTES:
                raise serializers.ValidationError({"file": "Expanded document is too large."})
    except BadZipFile:
        raise serializers.ValidationError({"file": "Document is invalid or corrupted."}) from None
    finally:
        upload.seek(0)


def _docx_lines(upload):
    document = Document(upload)
    yield from (paragraph.text for paragraph in document.paragraphs)
    for table in document.tables:
        for row in table.rows:
            yield " | ".join(cell.text for cell in row.cells)


def _xlsx_lines(upload):
    workbook = load_workbook(upload, read_only=True, data_only=True)
    try:
        for sheet in workbook.worksheets:
            yield f"[{sheet.title}]"
            for row in sheet.iter_rows(values_only=True):
                yield " | ".join("" if value is None else str(value) for value in row)
    finally:
        workbook.close()


def _pptx_lines(upload):
    presentation = Presentation(upload)
    for index, slide in enumerate(presentation.slides, start=1):
        yield f"[Slide {index}]"
        yield from (shape.text for shape in slide.shapes if hasattr(shape, "text"))


def extract_context_document(upload):
    extension = Path(upload.name).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise serializers.ValidationError({"file": "Use PDF, DOCX, XLSX, PPTX, TXT, Markdown, or CSV."})
    if upload.size > MAX_UPLOAD_BYTES:
        raise serializers.ValidationError({"file": "Document must not exceed 10 MB."})

    try:
        if extension in {".txt", ".md", ".csv"}:
            text, truncated = _bounded_text(upload.read().decode("utf-8-sig", errors="replace").splitlines())
        elif extension == ".pdf":
            reader = PdfReader(upload)
            if reader.is_encrypted:
                raise serializers.ValidationError({"file": "Password-protected PDF is not supported."})
            text, truncated = _bounded_text(page.extract_text() for page in reader.pages[:200])
        else:
            _validate_office_archive(upload)
            lines = {".docx": _docx_lines, ".xlsx": _xlsx_lines, ".pptx": _pptx_lines}[extension](upload)
            text, truncated = _bounded_text(lines)
    except serializers.ValidationError:
        raise
    except Exception:
        raise serializers.ValidationError({"file": "Document could not be read."}) from None
    if not text:
        raise serializers.ValidationError({"file": "Document does not contain readable text."})
    return {"name": Path(upload.name).name, "text": text, "truncated": truncated}
